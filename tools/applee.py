import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import plotly.express as px
from PIL import Image, ImageDraw, ImageFont


# Function to create a text animation GIF
def create_text_animation(texts, output_filename, duration=500):
    frames = []
    for text in texts:
        image = Image.new("RGB", (300, 100), "white")
        draw = ImageDraw.Draw(image)
        font = ImageFont.truetype("arial.ttf", 24)
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_position = ((300 - text_bbox[2]) // 2, (100 - text_bbox[3]) // 2)
        draw.text(text_position, text, font=font, fill="black")
        frames.append(image)

    frames[0].save(output_filename, save_all=True, append_images=frames[1:], duration=duration, loop=0)


# Function to create a content slide
def create_content_slide(prs, title, content):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content[0]

    for item in content[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    return slide


# Create a new presentation
prs = Presentation()

# Add title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Animated PowerPoint Presentation"
subtitle.text = "Created with Python"

# Create an animated chart using Plotly
df = px.data.gapminder()
fig = px.scatter(df, x="gdpPercap", y="lifeExp", animation_frame="year",
                 animation_group="country", size="pop", color="continent",
                 hover_name="country", log_x=True, size_max=55,
                 range_x=[100, 100000], range_y=[25, 90])

# Save the chart as an HTML file
fig.write_html("animated_chart.html")

# Create content slides
create_content_slide(prs, "Process Overview", ["Main Steps:", "Data Collection", "Analysis", "Visualization"])
create_content_slide(prs, "Data Collection", ["Methods:", "Surveys", "Interviews", "Observations"])
create_content_slide(prs, "Analysis", ["Techniques:", "Statistical Analysis", "Machine Learning", "Data Mining"])

# Create a slide for the animated chart
slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.title
title.text = "Global Development Trends"
subtitle = slide.shapes.subtitle
subtitle.text = "View the interactive chart at: [Your GitHub Pages URL]/animated_chart.html"

# Create text animation
texts = ["Step 1: Collect Data", "Step 2: Analyze", "Step 3: Visualize", "Process Complete!"]
create_text_animation(texts, "process_steps.gif")

# Add a slide with the text animation
slide = prs.slides.add_slide(prs.slide_layouts[6])
title = slide.shapes.title
title.text = "Process Steps"
left = top = Inches(1)
pic = slide.shapes.add_picture("process_steps.gif", left, top, height=Inches(5))

# Add a summary slide
summary_slide = create_content_slide(prs, "Summary",
                                     ["Key Takeaways:",
                                      "Understanding data is crucial",
                                      "Proper analysis leads to insights",
                                      "Visualization helps in communication",
                                      "Continuous learning is important"])

# Add a shape to emphasize a point
oval_shape = summary_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(5), Inches(1), Inches(0.5))
oval_shape.fill.solid()
oval_shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# Save the presentation
prs.save('animated_presentation.pptx')

print("Presentation created successfully!")